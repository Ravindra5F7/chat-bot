# # --- START OF FILE ai_core.py ---

# import os
# import logging
# import json
# import ollama
# import concurrent.futures
# from tqdm import tqdm
# from pypdf import PdfReader
# from pptx import Presentation
# from langchain_community.vectorstores import FAISS
# from langchain_ollama import OllamaEmbeddings, ChatOllama
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.docstore.document import Document
# from langchain.chains import LLMChain
# # PromptTemplate is already imported in config, but good practice if used directly here too
# # from langchain.prompts import PromptTemplate

# from config import (
#     OLLAMA_BASE_URL, OLLAMA_MODEL, OLLAMA_EMBED_MODEL, FAISS_FOLDER,
#     DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, RAG_CHUNK_K, MULTI_QUERY_COUNT,
#     ANALYSIS_MAX_CONTEXT_LENGTH, OLLAMA_REQUEST_TIMEOUT, RAG_SEARCH_K_PER_QUERY,
#     SUB_QUERY_PROMPT_TEMPLATE, SYNTHESIS_PROMPT_TEMPLATE, ANALYSIS_PROMPTS,
#     KG_CHUNK_SIZE, KG_CHUNK_OVERLAP, KG_MAX_WORKERS, KG_MODEL,
#     KG_OUTPUT_FOLDER, KG_PROMPT_TEMPLATE, KG_FILENAME_SUFFIX # Added KG_FILENAME_SUFFIX
# )
# from utils import parse_llm_response, escape_html # Assuming you have these

# logger = logging.getLogger(__name__)

# # Global State
# document_texts_cache = {}
# vector_store = None
# embeddings = None
# llm = None
# _kg_ollama_client = None

# # Initialization Functions
# def initialize_ai_components():
#     """Initializes embeddings and LLM."""
#     global embeddings, llm
#     logger.info("Initializing AI components...")
#     try:
#         embeddings = OllamaEmbeddings(
#             base_url=OLLAMA_BASE_URL,
#             model=OLLAMA_EMBED_MODEL,
#         )
#         llm = ChatOllama(
#             base_url=OLLAMA_BASE_URL,
#             model=OLLAMA_MODEL,
#             request_timeout=OLLAMA_REQUEST_TIMEOUT
#         )
#         logger.info("AI components initialized.")
#         return embeddings, llm
#     except Exception as e:
#         logger.error(f"Failed to initialize AI components: {e}", exc_info=True)
#         return None, None

# def load_vector_store():
#     """Loads or initializes FAISS vector store."""
#     global vector_store
#     try:
#         if os.path.exists(FAISS_FOLDER) and os.listdir(FAISS_FOLDER): # Check if directory is not empty
#             vector_store = FAISS.load_local(FAISS_FOLDER, embeddings, allow_dangerous_deserialization=True)
#             logger.info(f"Loaded existing FAISS vector store from {FAISS_FOLDER}.")
#         else:
#             logger.info(f"No FAISS index found at {FAISS_FOLDER} or it's empty. Starting with empty index.")
#             # Optionally initialize an empty store if needed immediately,
#             # or let it be created on first add_documents_to_vector_store
#         return True
#     except Exception as e:
#         logger.error(f"Failed to load FAISS vector store: {e}", exc_info=True)
#         return False

# def save_vector_store():
#     """Saves FAISS vector store."""
#     try:
#         if vector_store:
#             vector_store.save_local(FAISS_FOLDER)
#             logger.info(f"Saved FAISS vector store to {FAISS_FOLDER}.")
#     except Exception as e:
#         logger.error(f"Failed to save FAISS vector store: {e}", exc_info=True)

# def load_all_document_texts():
#     """Loads text from all documents into cache."""
#     import config # Explicit import
#     logger.info("Loading all document texts into cache...")
#     loaded_count = 0
#     for folder_path in [config.DEFAULT_PDFS_FOLDER, config.UPLOAD_FOLDER]:
#         if not os.path.exists(folder_path):
#             logger.warning(f"Document folder not found: {folder_path}")
#             continue
#         for filename in os.listdir(folder_path):
#             if filename.lower().endswith(tuple(config.ALLOWED_EXTENSIONS)) and not filename.startswith('~'):
#                 filepath = os.path.join(folder_path, filename)
#                 if filename not in document_texts_cache: # Avoid reloading if already cached
#                     text = extract_text_from_file(filepath) # Renamed for clarity
#                     if text:
#                         document_texts_cache[filename] = text
#                         loaded_count +=1
#     logger.info(f"Loaded texts for {loaded_count} new documents into cache. Total cached: {len(document_texts_cache)}")


# # Text Extraction and Chunking
# def extract_text_from_file(filepath: str) -> str: # Renamed from extract_text_from_pdf
#     """Extracts text from PDF or PPT file."""
#     try:
#         ext = os.path.splitext(filepath)[1].lower()
#         if ext == '.pdf':
#             reader = PdfReader(filepath)
#             if not reader.pages:
#                 logger.error(f"PDF file {filepath} has no pages or could not be read.")
#                 return ""
#             text = ""
#             for i, page in enumerate(reader.pages):
#                 page_text = page.extract_text()
#                 if page_text:
#                     text += page_text + "\n"
#                 else:
#                     logger.warning(f"No text extracted from page {i+1} of {filepath}.")
#             return text.strip() or ""
#         elif ext in ('.pptx', '.ppt'):
#             prs = Presentation(filepath)
#             text = ""
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         text += shape.text + "\n"
#             return text.strip() or ""
#         else:
#             logger.warning(f"Unsupported file extension for text extraction: {ext} in {filepath}")
#             return ""
#     except FileNotFoundError:
#         logger.error(f"File not found for text extraction: {filepath}")
#         return ""
#     except Exception as e:
#         logger.error(f"Error extracting text from '{filepath}': {e}", exc_info=True)
#         return ""

# def create_chunks_from_text(text: str, filename: str) -> list[Document]:
#     """Creates document chunks from text."""
#     try:
#         # TODO: Consider making chunk_size and chunk_overlap configurable from config.py for RAG
#         splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#         chunks = splitter.split_text(text)
#         return [Document(page_content=chunk, metadata={"source": filename}) for chunk in chunks]
#     except Exception as e:
#         logger.error(f"Error creating chunks for '{filename}': {e}", exc_info=True)
#         return []

# def add_documents_to_vector_store(documents: list[Document]) -> bool:
#     """Adds documents to FAISS vector store."""
#     global vector_store
#     if not embeddings:
#         logger.error("Cannot add documents to vector store: Embeddings not initialized.")
#         return False
#     try:
#         if not vector_store:
#             logger.info("Creating new FAISS vector store.")
#             vector_store = FAISS.from_documents(documents, embeddings)
#         else:
#             vector_store.add_documents(documents)
#         save_vector_store()
#         return True
#     except Exception as e:
#         logger.error(f"Error adding documents to vector store: {e}", exc_info=True)
#         return False

# # RAG and Analysis
# def generate_sub_queries(query: str) -> list[str]:
#     """Generates sub-queries for RAG."""
#     if not llm:
#         logger.error("LLM not initialized, cannot generate sub-queries.")
#         return []
#     try:
#         chain = LLMChain(llm=llm, prompt=SUB_QUERY_PROMPT_TEMPLATE)
#         response = chain.run(query=query, num_queries=MULTI_QUERY_COUNT)
#         return [q.strip() for q in response.split('\n') if q.strip()]
#     except Exception as e:
#         logger.error(f"Error generating sub-queries: {e}", exc_info=True)
#         return []

# def perform_rag_search(query: str) -> tuple[list[Document], str, dict]:
#     """Performs RAG search."""
#     if not vector_store:
#         logger.warning("Vector store not loaded/initialized. RAG search cannot be performed.")
#         return [], "Knowledge base is currently unavailable.", {}
#     try:
#         sub_queries = generate_sub_queries(query) if MULTI_QUERY_COUNT > 0 else []
#         all_queries = [query] + sub_queries
        
#         all_retrieved_docs_with_scores = []
#         for q_idx, q_text in enumerate(all_queries):
#             # Using similarity_search_with_score to potentially de-duplicate or rank later
#             results_with_scores = vector_store.similarity_search_with_score(q_text, k=RAG_SEARCH_K_PER_QUERY)
#             all_retrieved_docs_with_scores.extend(results_with_scores)

#         # De-duplicate documents based on content and source, keeping the one with the best score (lower is better for L2/cosine)
#         unique_docs_dict = {}
#         for doc, score in all_retrieved_docs_with_scores:
#             doc_key = (doc.metadata['source'], doc.page_content)
#             if doc_key not in unique_docs_dict or score < unique_docs_dict[doc_key][1]:
#                 unique_docs_dict[doc_key] = (doc, score)
        
#         # Sort unique documents by score and take top RAG_CHUNK_K
#         sorted_unique_docs = sorted(list(unique_docs_dict.values()), key=lambda item: item[1])
#         final_docs = [item[0] for item in sorted_unique_docs[:RAG_CHUNK_K]]

#         if not final_docs:
#             logger.info(f"No relevant documents found for query: {query}")
#             return [], "No specific documents found for your query.", {}

#         context_text = "\n\n".join(f"[Chunk {i+1} from: {doc.metadata['source']}]\n{doc.page_content}" for i, doc in enumerate(final_docs))
#         # Map for citation: Key is citation number (1-based), value is details
#         docs_map = {i+1: {"source": doc.metadata['source'], "content": doc.page_content} for i, doc in enumerate(final_docs)}
        
#         return final_docs, context_text, docs_map
#     except Exception as e:
#         logger.error(f"Error in RAG search: {e}", exc_info=True)
#         return [], "Error retrieving documents.", {}

# # --- KG Related Functions ---
# def get_kg_filepath(doc_filename: str) -> str:
#     """Gets the expected filepath for a document's KG."""
#     base_name = os.path.splitext(doc_filename)[0]
#     kg_file = f"{base_name}{KG_FILENAME_SUFFIX}"
#     return os.path.join(KG_OUTPUT_FOLDER, kg_file)

# def load_knowledge_graph(doc_filename: str) -> dict | None:
#     """Loads the knowledge graph for a given document filename."""
#     kg_filepath = get_kg_filepath(doc_filename)
#     if os.path.exists(kg_filepath):
#         try:
#             with open(kg_filepath, 'r', encoding='utf-8') as f:
#                 kg_data = json.load(f)
#             logger.info(f"Loaded KG for '{doc_filename}' from {kg_filepath}")
#             return kg_data
#         except Exception as e:
#             logger.error(f"Error loading KG JSON from {kg_filepath}: {e}", exc_info=True)
#     else:
#         logger.debug(f"No KG found for '{doc_filename}' at {kg_filepath}")
#     return None

# def format_kg_info_for_llm(kg_data: dict, query: str, max_insights: int = 3, max_details_per_node: int = 2) -> str:
#     """
#     Extracts and formats relevant KG information based on the query for LLM context.
#     """
#     if not kg_data or not isinstance(kg_data, dict) or 'nodes' not in kg_data:
#         return ""

#     query_keywords = set(q.lower() for q in query.split() if len(q) > 3) # Keywords from query
    
#     relevant_node_infos = []
#     nodes_by_id = {node.get('id'): node for node in kg_data.get('nodes', []) if isinstance(node, dict) and node.get('id')}

#     # Find nodes relevant to the query
#     for node_id, node_data in nodes_by_id.items():
#         description = node_data.get('description', '').lower()
#         node_id_lower = str(node_id).lower()
#         if any(keyword in description for keyword in query_keywords) or \
#            any(keyword in node_id_lower for keyword in query_keywords):
#             relevant_node_infos.append(node_data)
    
#     if not relevant_node_infos:
#         return ""

#     output_str = "Key concepts and relationships from Knowledge Graph:\n"
#     insights_count = 0

#     for node in relevant_node_infos:
#         if insights_count >= max_insights:
#             break
        
#         node_id = node.get('id')
#         desc = node.get('description', 'N/A')
#         node_type = node.get('type', 'Concept')
#         output_str += f"- Node: '{node_id}' (Type: {node_type})\n  Description: {desc}\n"
        
#         # Find connected edges
#         related_count = 0
#         for edge in kg_data.get('edges', []):
#             if related_count >= max_details_per_node: break
#             if not isinstance(edge, dict): continue
            
#             rel_from = edge.get('from')
#             rel_to = edge.get('to')
#             relationship = edge.get('relationship')

#             if rel_from == node_id and rel_to in nodes_by_id:
#                 output_str += f"  - Relates to '{rel_to}' (via: {relationship})\n"
#                 related_count +=1
#             elif rel_to == node_id and rel_from in nodes_by_id:
#                 output_str += f"  - Is related from '{rel_from}' (via: {relationship})\n"
#                 related_count +=1
#         insights_count += 1
    
#     return output_str.strip() if insights_count > 0 else ""

# # --- End KG Related Functions ---


# def synthesize_chat_response(query: str, chat_history: list, tool_output: str = "") -> tuple[str, str]:
#     """Synthesizes chat response using RAG context and KG insights, or proposes actions."""
#     if not llm:
#         logger.error("LLM not initialized, cannot synthesize response.")
#         return "Error: AI model is not available.", ""
#     try:
#         # The prompt template now handles all context (chat_history, tool_output)
#         # Remove old context building logic (kg_derived_insights, enriched_context)
        
#         chain = LLMChain(llm=llm, prompt=SYNTHESIS_PROMPT_TEMPLATE)
#         # Pass the necessary variables to the prompt template
#         response_text = chain.run(
#             query=query,
#             chat_history=chat_history,
#             tool_output=tool_output
#         ) 
        
#         # Parse LLM's ReAct output
#         thought = ""
#         action = ""
#         speak = ""

#         lines = response_text.split('\n')
#         current_section = None
#         for line in lines:
#             if line.strip().startswith("Thought:"):
#                 current_section = "thought"
#                 thought += line.strip()[len("Thought:"):].strip() + "\n"
#             elif line.strip().startswith("Action:"):
#                 current_section = "action"
#                 action += line.strip()[len("Action:"):].strip() + "\n"
#             elif line.strip().startswith("Speak:"):
#                 current_section = "speak"
#                 speak += line.strip()[len("Speak:"):].strip() + "\n"
#             elif current_section == "thought":
#                 thought += line.strip() + "\n"
#             elif current_section == "action":
#                 action += line.strip() + "\n"
#             elif current_section == "speak":
#                 speak += line.strip() + "\n"
        
#         # Clean up trailing newlines
#         thought = thought.strip()
#         action = action.strip()
#         speak = speak.strip()

#         # Return the parsed components. The calling service will decide what to do.
#         # We'll use a special format for the returned string to indicate the action/speak.
#         # The 'thinking' part will be the thought.
#         if action: # If an action is proposed, return it for execution
#             # We'll use a specific prefix to signal an action to the Flask app
#             return f"<ACTION>{action}</ACTION>", thought
#         elif speak: # If the LLM wants to speak, return the final answer
#             return speak, thought
#         else: # Fallback for unexpected output
#             logger.warning(f"LLM returned unparseable response: {response_text[:200]}...")
#             return "Sorry, I could not understand the response from the AI. Please try again.", thought or "No clear thought or action/speak found."

#     except Exception as e:
#         logger.error(f"Error synthesizing chat response: {e}", exc_info=True)
#         return f"Error: Could not generate a response. {str(e)}", ""


# def generate_document_analysis(filename: str, analysis_type: str) -> tuple[str, str]:
#     """Generates document analysis."""
#     import config # Explicit import
#     if not llm:
#         logger.error("LLM not initialized, cannot generate document analysis.")
#         return "Error: AI model is not available.", ""
#     try:
#         text = document_texts_cache.get(filename)
#         if not text:
#             # Attempt to load from disk if not in cache
#             file_path_to_try = os.path.join(UPLOAD_FOLDER, filename)
#             if not os.path.exists(file_path_to_try):
#                 file_path_to_try = os.path.join(DEFAULT_PDFS_FOLDER, filename)
            
#             if os.path.exists(file_path_to_try):
#                 logger.info(f"Text for '{filename}' not in cache, loading from {file_path_to_try} for analysis.")
#                 text = extract_text_from_file(file_path_to_try)
#                 if text:
#                     document_texts_cache[filename] = text # Cache it now
#             else:
#                 logger.error(f"File '{filename}' not found for analysis.")
#                 return None, "File not found."

#         if not text:
#             return None, "No text available for analysis."

#         chain = LLMChain(llm=llm, prompt=config.ANALYSIS_PROMPTS[analysis_type])
#         # Ensure text is not excessively long for this type of analysis
#         response_text = chain.run(doc_text_for_llm=text[:ANALYSIS_MAX_CONTEXT_LENGTH]) 
#         thinking = response_text.split('</thinking>')[0].replace('<thinking>', '').strip() if '</thinking>' in response_text else ""
#         content = response_text.split('</thinking>')[-1].strip()
#         return content, thinking
#     except Exception as e:
#         logger.error(f"Error in document analysis for '{filename}': {e}", exc_info=True)
#         return f"Error: {str(e)}", ""

# # Knowledge Graph Generation (largely same, but saving part needs update)
# def _initialize_kg_ollama_client() -> ollama.Client | None:
#     global _kg_ollama_client
#     if _kg_ollama_client:
#         return _kg_ollama_client
#     try:
#         logger.info(f"Initializing KG Ollama client: {OLLAMA_BASE_URL}, model={KG_MODEL}")
#         _kg_ollama_client = ollama.Client(host=OLLAMA_BASE_URL)
#         _kg_ollama_client.list() # Test connection
#         logger.info("KG Ollama client initialized.")
#         return _kg_ollama_client
#     except Exception as e:
#         logger.error(f"Failed to initialize KG Ollama client: {e}", exc_info=True)
#         _kg_ollama_client = None
#         return None

# def _kg_split_into_chunks(text: str, chunk_size: int, overlap: int) -> list[str]:
#     logger.info(f"KG: Splitting text into chunks (size={chunk_size}, overlap={overlap})...")
#     chunks = []
#     start = 0
#     text_len = len(text)
#     while start < text_len:
#         end = min(start + chunk_size, text_len)
#         chunks.append(text[start:end])
#         next_start = end - overlap
#         # Ensure progress, especially with small texts or large overlaps
#         if next_start <= start and end < text_len :
#             start = start + 1 
#         else:
#             start = next_start if end < text_len else end
#     logger.info(f"KG: Split into {len(chunks)} chunks.")
#     return chunks


# def _kg_process_single_chunk(chunk_data: tuple[int, str], ollama_client_instance: ollama.Client, model_name: str, prompt_template_str: str) -> dict | None:
#     index, chunk_text = chunk_data
#     chunk_num = index + 1
    
#     # The KG_PROMPT_TEMPLATE from config.py is already escaped correctly
#     full_prompt = prompt_template_str.format(chunk_text=chunk_text) 
    
#     if not ollama_client_instance:
#         logger.error(f"KG: Ollama client not available for chunk {chunk_num}.")
#         return None
#     try:
#         response = ollama_client_instance.chat(
#             model=model_name,
#             messages=[{"role": "user", "content": full_prompt}],
#             format="json", # Ollama should handle this based on model capabilities
#             options={"num_ctx": 4096, "temperature": 0.3} # Example options
#         )
#         content = response.get('message', {}).get('content', '')
#         if not content:
#             logger.warning(f"KG: Empty response for chunk {chunk_num}")
#             return {}
        
#         # Basic cleanup for ```json ... ``` markdown blocks if model adds them
#         if content.strip().startswith('```json'):
#             content = content.strip()[7:-3].strip()
#         elif content.strip().startswith('```'): # More generic ``` removal
#             content = content.strip()[3:-3].strip()
            
#         graph_data = json.loads(content) # Expecting JSON directly now
        
#         # Validate basic structure
#         if isinstance(graph_data, dict) and \
#            'nodes' in graph_data and isinstance(graph_data['nodes'], list) and \
#            'edges' in graph_data and isinstance(graph_data['edges'], list):
#             return graph_data
        
#         logger.warning(f"KG: Invalid graph structure for chunk {chunk_num}. Content: {content[:200]}...")
#         return {}
#     except json.JSONDecodeError as je:
#         logger.error(f"KG: JSON Decode Error processing chunk {chunk_num}: {je}. Content: {content[:500]}...", exc_info=False) # Don't need full exc_info for common JSON error
#         return {}
#     except Exception as e:
#         logger.error(f"KG: Error processing chunk {chunk_num}: {e}", exc_info=True)
#         return {}

# def _kg_merge_graphs(graphs: list[dict]) -> dict:
#     logger.info("KG: Merging graph fragments...")
#     final_nodes = {} # Use dict for easy de-duplication and update
#     final_edges = set() # Use set of tuples for de-duplication

#     for i, graph_fragment in enumerate(graphs):
#         if graph_fragment is None:
#             logger.warning(f"KG: Skipping None graph fragment at index {i}.")
#             continue
#         if not isinstance(graph_fragment, dict) or 'nodes' not in graph_fragment or 'edges' not in graph_fragment:
#             logger.warning(f"KG: Skipping invalid graph fragment at index {i}")
#             continue

#         # Process nodes
#         for node in graph_fragment.get('nodes', []):
#             if not isinstance(node, dict):
#                 logger.warning(f"KG: Skipping non-dict node in fragment {i}: {node}")
#                 continue
#             node_id = node.get('id')
#             if not node_id or not isinstance(node_id, str): # Ensure ID is a non-empty string
#                 logger.warning(f"KG: Skipping node with invalid or missing ID in fragment {i}: {node}")
#                 continue
            
#             if node_id not in final_nodes:
#                 final_nodes[node_id] = node
#             else:
#                 # Merge properties, e.g., longer description, fill missing parent/type
#                 existing_node = final_nodes[node_id]
#                 for key in ['description', 'type', 'parent']:
#                     if node.get(key) and (not existing_node.get(key) or len(str(node.get(key))) > len(str(existing_node.get(key)))):
#                          if key == 'parent' and node.get(key) == node_id: continue # Avoid self-parenting from merge
#                          existing_node[key] = node.get(key)


#         # Process edges
#         for edge in graph_fragment.get('edges', []):
#             if not isinstance(edge, dict) or not all(k in edge for k in ['from', 'to', 'relationship']):
#                 logger.warning(f"KG: Skipping invalid edge (missing keys) in fragment {i}: {edge}")
#                 continue
#             # Ensure all parts of an edge are strings and not empty
#             if not all(isinstance(edge.get(k), str) and edge.get(k) for k in ['from', 'to', 'relationship']):
#                 logger.warning(f"KG: Skipping edge with non-string or empty components in fragment {i}: {edge}")
#                 continue
#             # Avoid self-loops if not desired, or ensure nodes exist
#             if edge['from'] == edge['to']: # Optional: disallow self-loops
#                 # logger.debug(f"KG: Skipping self-loop edge: {edge}")
#                 continue
#             edge_tuple = (edge['from'], edge['to'], edge['relationship'])
#             final_edges.add(edge_tuple)

#     # Convert back to list format
#     merged_graph_nodes = list(final_nodes.values())
#     merged_graph_edges = [{"from": e[0], "to": e[1], "relationship": e[2]} for e in final_edges]
    
#     # Filter out edges where 'from' or 'to' node doesn't exist in the final_nodes list
#     # This can happen if chunks produce partial relationships
#     valid_node_ids = set(n['id'] for n in merged_graph_nodes)
#     final_merged_edges = [edge for edge in merged_graph_edges if edge['from'] in valid_node_ids and edge['to'] in valid_node_ids]

#     merged_graph = {"nodes": merged_graph_nodes, "edges": final_merged_edges}
#     logger.info(f"KG: Merged into {len(merged_graph['nodes'])} nodes, {len(merged_graph['edges'])} edges.")
#     return merged_graph


# def _kg_save_graph(graph: dict, original_doc_filename: str):
#     """Saves the graph to a file named after the original document."""
#     kg_filepath = get_kg_filepath(original_doc_filename)
#     try:
#         logger.info(f"KG: Saving graph for '{original_doc_filename}' to {kg_filepath}...")
#         os.makedirs(os.path.dirname(kg_filepath), exist_ok=True)
#         with open(kg_filepath, 'w', encoding='utf-8') as f:
#             json.dump(graph, f, indent=2, ensure_ascii=False)
#         logger.info(f"KG: Graph for '{original_doc_filename}' saved to {kg_filepath}.")
#     except Exception as e:
#         logger.error(f"Failed to save graph to {kg_filepath}: {e}", exc_info=True)


# def generate_knowledge_graph_from_pdf(doc_filename: str) -> dict | None: # Renamed from generate_knowledge_graph_from_pdf for clarity
#     """
#     Generates a knowledge graph from a PDF or PPT file, saving it named after the doc.
#     'doc_filename' is the base name of the file, e.g., "MyDocument.pdf".
#     """
#     global document_texts_cache
#     kg_ollama_client = _initialize_kg_ollama_client()
#     if not kg_ollama_client:
#         logger.error(f"KG: Cannot generate graph for '{doc_filename}'. KG Ollama client not initialized.")
#         return None

#     logger.info(f"KG: Generating Knowledge Graph for '{doc_filename}'...")
    
#     doc_text = document_texts_cache.get(doc_filename)
#     if not doc_text:
#         logger.debug(f"KG: '{doc_filename}' not in cache. Attempting to load from disk...")
#         # Determine full path (check UPLOAD_FOLDER first, then DEFAULT_PDFS_FOLDER)
#         file_path_to_load = os.path.join(UPLOAD_FOLDER, doc_filename)
#         if not os.path.exists(file_path_to_load):
#             file_path_to_load = os.path.join(DEFAULT_PDFS_FOLDER, doc_filename)

#         if os.path.exists(file_path_to_load):
#             logger.info(f"KG: Found '{doc_filename}' at {file_path_to_load}. Extracting text...")
#             doc_text = extract_text_from_file(file_path_to_load)
#             if doc_text:
#                 document_texts_cache[doc_filename] = doc_text # Cache it
#                 logger.info(f"KG: Cached text for '{doc_filename}' from {file_path_to_load}.")
#             else:
#                 logger.error(f"KG: Failed to extract text from '{doc_filename}' at {file_path_to_load}.")
#                 return None
#         else:
#             logger.error(f"KG: File '{doc_filename}' not found in {UPLOAD_FOLDER} or {DEFAULT_PDFS_FOLDER}.")
#             return None

#     if not doc_text.strip(): # Check if text is not just whitespace
#         logger.error(f"KG: No actual text content available for '{doc_filename}' to generate graph.")
#         return None

#     chunks = _kg_split_into_chunks(doc_text, chunk_size=KG_CHUNK_SIZE, overlap=KG_CHUNK_OVERLAP)
#     if not chunks:
#         logger.warning(f"KG: No chunks generated for '{doc_filename}'. KG generation aborted.")
#         return None

#     logger.info(f"KG: Processing {len(chunks)} chunks for '{doc_filename}' using {KG_MAX_WORKERS} workers with model {KG_MODEL}.")
#     all_partial_graphs = []
#     # Use KG_PROMPT_TEMPLATE from config
#     prompt_for_chunks = KG_PROMPT_TEMPLATE 

#     with concurrent.futures.ThreadPoolExecutor(max_workers=KG_MAX_WORKERS) as executor:
#         # Prepare tasks: tuple of (index, chunk_text)
#         tasks = [(i, chunk) for i, chunk in enumerate(chunks)]
#         # Submit tasks and map futures back to their original index for logging/debugging
#         future_to_index = {
#             executor.submit(_kg_process_single_chunk, task_data, kg_ollama_client, KG_MODEL, prompt_for_chunks): task_data[0] 
#             for task_data in tasks
#         }
        
#         for future in tqdm(concurrent.futures.as_completed(future_to_index), total=len(tasks), desc=f"KG: Processing Chunks ({doc_filename})", unit="chunk"):
#             original_chunk_index = future_to_index[future]
#             try:
#                 graph_data_fragment = future.result()
#                 if graph_data_fragment:
#                     all_partial_graphs.append(graph_data_fragment)
#             except Exception as e: # This catches exceptions from _kg_process_single_chunk if not caught inside
#                 logger.error(f"KG: Uncaught error from chunk {original_chunk_index + 1} processing for '{doc_filename}': {e}", exc_info=True)

#     logger.info(f"KG: Processed {len(all_partial_graphs)} valid graph fragments for '{doc_filename}' from {len(chunks)} chunks.")

#     if not all_partial_graphs:
#         logger.error(f"KG: No valid graph fragments generated for '{doc_filename}'. Cannot merge.")
#         return None

#     final_graph = _kg_merge_graphs(all_partial_graphs)
    
#     # Save graph using the original document filename to create the KG filename
#     _kg_save_graph(final_graph, doc_filename) 

#     logger.info(f"KG: Knowledge graph generation for '{doc_filename}' complete.")
#     return final_graph


# # Main Execution / Test
# if __name__ == "__main__":
#     # Basic logging for testing
#     logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - [%(name)s:%(lineno)d] - %(message)s")
    
#     # Setup directories (idempotent)
#     for dir_path in [DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, FAISS_FOLDER, KG_OUTPUT_FOLDER]:
#         os.makedirs(dir_path, exist_ok=True)
    
#     logger.info("AI Core Testing (with KG Integration)...")
    
#     # 1. Initialize AI components (LLM, Embeddings)
#     embeddings_instance, llm_instance = initialize_ai_components()
#     if not embeddings_instance or not llm_instance:
#         logger.error("Failed to initialize core AI components. Aborting test.")
#         exit()
    
#     # 2. Load vector store (it might be empty initially)
#     load_vector_store()

#     # 3. Create a dummy PDF for testing (if it doesn't exist)
#     dummy_pdf_name = "DevOps-Test-KG.pdf"
#     dummy_pdf_path = os.path.join(DEFAULT_PDFS_FOLDER, dummy_pdf_name)
    
#     # if not os.path.exists(dummy_pdf_path):
#     #     try:
#     #         from reportlab.pdfgen import canvas
#     #         from reportlab.lib.pagesizes import letter
#     #         logger.info(f"Creating dummy PDF: {dummy_pdf_path}")
#     #         c = canvas.Canvas(dummy_pdf_path, pagesize=letter)
#     #         c.drawString(72, 800, "DevOps Test Document for Knowledge Graph")
#     #         c.drawString(72, 780, "This document covers Continuous Integration (CI) and Continuous Deployment (CD).")
#     #         c.drawString(72, 760, "Key tools include Jenkins for CI, and Docker for containerization.")
#     #         c.drawString(72, 740, "Kubernetes is used for orchestrating Docker containers.")
#     #         c.drawString(72, 720, "Monitoring is crucial in DevOps, often using Prometheus.")
#     #         c.save()
#     #         logger.info(f"Dummy PDF '{dummy_pdf_name}' created.")
#     #     except ImportError:
#     #         logger.warning("ReportLab not installed. Cannot create dummy PDF. Please create it manually for testing.")
#     #     except Exception as e:
#     #         logger.error(f"Error creating dummy PDF: {e}")

#     if os.path.exists(dummy_pdf_path):
#         # 4. Process the dummy PDF: extract text, chunk, add to vector store
#         logger.info(f"Processing '{dummy_pdf_name}' for RAG...")
#         text_content = extract_text_from_file(dummy_pdf_path)
#         if text_content:
#             document_texts_cache[dummy_pdf_name] = text_content # Cache it
#             chunks = create_chunks_from_text(text_content, dummy_pdf_name)
#             if chunks:
#                 add_documents_to_vector_store(chunks)
#                 logger.info(f"'{dummy_pdf_name}' processed and added to vector store.")
#             else:
#                 logger.error(f"Could not create chunks for '{dummy_pdf_name}'.")
#         else:
#             logger.error(f"Could not extract text from '{dummy_pdf_name}'.")

#         # 5. Generate Knowledge Graph for the dummy PDF
#         logger.info(f"Testing KG generation for {dummy_pdf_name}...")
#         kg_result = generate_knowledge_graph_from_pdf(dummy_pdf_name) # Pass base filename
#         if kg_result:
#             logger.info(f"KG for '{dummy_pdf_name}' generated: {len(kg_result.get('nodes',[]))} nodes, {len(kg_result.get('edges',[]))} edges.")
#             kg_file_expected = get_kg_filepath(dummy_pdf_name)
#             logger.info(f"KG saved to: {kg_file_expected}")
#         else:
#             logger.error(f"Failed to generate KG for '{dummy_pdf_name}'.")
#     else:
#         logger.warning(f"Dummy PDF '{dummy_pdf_name}' not found. Skipping RAG and KG generation tests for it.")

#     # 6. Test chat response with RAG and KG integration
#     test_query = "What is CI/CD and what tools are used?"
#     logger.info(f"\nTesting chat with query: '{test_query}'")
    
#     if vector_store: # Ensure RAG can even attempt
#         rag_docs, rag_context_text, rag_docs_map = perform_rag_search(test_query)
#         logger.info(f"RAG found {len(rag_docs)} documents for the query.")
#         logger.debug(f"RAG context text: {rag_context_text[:300]}...")
        
#         answer, thinking = synthesize_chat_response(test_query, [], "")
        
#         logger.info(f"\n--- Test Chat Response ---")
#         logger.info(f"Thinking:\n{thinking}")
#         logger.info(f"Answer:\n{answer}")
#         logger.info(f"--- End Test Chat Response ---")
#     else:
#         logger.warning("Vector store not available. Skipping chat test.")
        
#     logger.info("\nAI Core test (with KG integration) completed.")

# # --- END OF FILE ai_core.py ---

# --- START OF FILE ai_core.py ---

import os
import logging
import json
import ollama # For KG client
import concurrent.futures
from tqdm import tqdm # For KG progress
from pypdf import PdfReader # For PDF text extraction
from pptx import Presentation # For PPTX text extraction

# Langchain components
from langchain_community.vectorstores import FAISS
from langchain_ollama import OllamaEmbeddings, ChatOllama # For main LLM and embeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate # Import if used directly here

# Import from your project's config
# Ensure these are correctly defined in your config.py
from config import (
    OLLAMA_BASE_URL, OLLAMA_MODEL, OLLAMA_EMBED_MODEL, FAISS_FOLDER,
    DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, RAG_CHUNK_K, MULTI_QUERY_COUNT,
    ANALYSIS_MAX_CONTEXT_LENGTH, OLLAMA_REQUEST_TIMEOUT, RAG_SEARCH_K_PER_QUERY,
    SUB_QUERY_PROMPT_TEMPLATE, SYNTHESIS_PROMPT_TEMPLATE, ANALYSIS_PROMPTS,
    KG_CHUNK_SIZE, KG_CHUNK_OVERLAP, KG_MAX_WORKERS, KG_MODEL, # Ensure KG_MODEL is defined
    KG_OUTPUT_FOLDER, KG_PROMPT_TEMPLATE, KG_FILENAME_SUFFIX, ALLOWED_EXTENSIONS
)
from utils import parse_llm_response, escape_html # Assuming these exist in utils.py

logger = logging.getLogger(__name__)

# Global State Variables
document_texts_cache = {} # Cache for extracted text from documents
vector_store = None       # FAISS vector store instance
embeddings = None         # Langchain Embeddings instance (OllamaEmbeddings)
llm = None                # Langchain LLM instance (ChatOllama for main chat)
_kg_ollama_client = None  # Separate Ollama client for KG generation

# --- Initialization Functions ---
def initialize_ai_components():
    global embeddings, llm
    logger.info("[ai_core] Initializing AI components (Embeddings, Main LLM)...")
    try:
        logger.info(f"[ai_core] Embedding Model: {OLLAMA_EMBED_MODEL}, Main LLM Model: {OLLAMA_MODEL}, Base URL: {OLLAMA_BASE_URL}")
        embeddings = OllamaEmbeddings(
            base_url=OLLAMA_BASE_URL,
            model=OLLAMA_EMBED_MODEL,
        )
        llm = ChatOllama( # Main LLM for chat synthesis and analysis
            base_url=OLLAMA_BASE_URL,
            model=OLLAMA_MODEL,
            request_timeout=OLLAMA_REQUEST_TIMEOUT, # Make sure this is defined in config
            # Add other parameters like temperature if needed, e.g., temperature=0.7
        )
        # Test connection by listing models or a small generation
        llm.invoke("Hello") 
        embeddings.embed_query("test")
        logger.info("[ai_core] Main LLM and Embeddings initialized and tested successfully.")
        return embeddings, llm
    except Exception as e:
        logger.error(f"[ai_core] Failed to initialize AI components: {e}", exc_info=True)
        embeddings = None
        llm = None
        return None, None

def _initialize_kg_ollama_client() -> ollama.Client | None:
    """Initializes a separate Ollama client specifically for KG generation if not already done."""
    global _kg_ollama_client
    if _kg_ollama_client:
        return _kg_ollama_client
    try:
        logger.info(f"[ai_core] Initializing KG Ollama client: Base URL={OLLAMA_BASE_URL}, Model={KG_MODEL}")
        if not KG_MODEL: # KG_MODEL must be defined in config.py
            logger.error("[ai_core] KG_MODEL is not defined in configuration. Cannot initialize KG client.")
            return None
        _kg_ollama_client = ollama.Client(host=OLLAMA_BASE_URL) # ollama.Client uses 'host'
        _kg_ollama_client.list() # Test connection
        logger.info("[ai_core] KG Ollama client initialized successfully.")
        return _kg_ollama_client
    except Exception as e:
        logger.error(f"[ai_core] Failed to initialize KG Ollama client: {e}", exc_info=True)
        _kg_ollama_client = None
        return None

def load_vector_store():
    global vector_store
    if not embeddings:
        logger.error("[ai_core] Cannot load vector store: Embeddings not initialized.")
        return False
    try:
        # Ensure FAISS_FOLDER path exists before trying to load
        os.makedirs(FAISS_FOLDER, exist_ok=True)
        
        faiss_index_file = os.path.join(FAISS_FOLDER, "index.faiss")
        faiss_pkl_file = os.path.join(FAISS_FOLDER, "index.pkl")

        if os.path.exists(faiss_index_file) and os.path.exists(faiss_pkl_file):
            vector_store = FAISS.load_local(
                FAISS_FOLDER, embeddings, allow_dangerous_deserialization=True
            )
            idx_size = vector_store.index.ntotal if vector_store.index else 0
            logger.info(f"[ai_core] Loaded existing FAISS vector store from {FAISS_FOLDER}. Index size: {idx_size}")
        else:
            logger.info(f"[ai_core] No complete FAISS index found at {FAISS_FOLDER}. Will start with an empty/new index if documents are added.")
            vector_store = None # Explicitly set to None if not loaded
        return True
    except Exception as e:
        logger.error(f"[ai_core] Failed to load FAISS vector store: {e}", exc_info=True)
        vector_store = None
        return False

def save_vector_store():
    if not vector_store:
        logger.warning("[ai_core] No vector store instance to save.")
        return
    if not embeddings: # Should not happen if vector_store exists
        logger.error("[ai_core] Cannot save vector store: Embeddings not initialized.")
        return
    try:
        os.makedirs(FAISS_FOLDER, exist_ok=True)
        vector_store.save_local(FAISS_FOLDER)
        logger.info(f"[ai_core] Saved FAISS vector store to {FAISS_FOLDER}.")
    except Exception as e:
        logger.error(f"[ai_core] Failed to save FAISS vector store: {e}", exc_info=True)

def load_all_document_texts():
    global document_texts_cache
    logger.info("[ai_core] Loading all document texts into cache...")
    loaded_count = 0
    newly_loaded_files = []
    for folder_path in [DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER]:
        if not folder_path or not os.path.exists(folder_path): # Check if path is configured and exists
            logger.warning(f"[ai_core] Document folder not found or not configured: {folder_path}")
            continue
        for filename in os.listdir(folder_path):
            # Use ALLOWED_EXTENSIONS from config.py (should be list of '.ext')
            if filename.lower().endswith(tuple(ALLOWED_EXTENSIONS)) and not filename.startswith('~'):
                filepath = os.path.join(folder_path, filename)
                if filename not in document_texts_cache:
                    text = extract_text_from_file(filepath)
                    if text:
                        document_texts_cache[filename] = text
                        loaded_count += 1
                        newly_loaded_files.append(filename)
    if newly_loaded_files:
        logger.info(f"[ai_core] Loaded texts for {loaded_count} new documents: {', '.join(newly_loaded_files)}")
    logger.info(f"[ai_core] Total documents in text cache: {len(document_texts_cache)}")

# --- Text Extraction and Chunking ---
def extract_text_from_file(filepath: str) -> str | None:
    logger.debug(f"[ai_core] Attempting to extract text from: {filepath}")
    try:
        ext = os.path.splitext(filepath)[1].lower()
        text_content = ""
        if ext == '.pdf':
            reader = PdfReader(filepath)
            if not reader.pages:
                logger.error(f"[ai_core] PDF file {filepath} has no pages or is corrupted.")
                return None
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text + "\n"
                else:
                    logger.debug(f"[ai_core] No text extracted from page {i+1} of {filepath}.")
            return text_content.strip() if text_content.strip() else None
        elif ext in ('.pptx', '.ppt'):
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            return text_content.strip() if text_content.strip() else None
        else:
            logger.warning(f"[ai_core] Unsupported file type for text extraction: {ext} for file {filepath}")
            return None
    except FileNotFoundError:
        logger.error(f"[ai_core] File not found for text extraction: {filepath}")
        return None
    except Exception as e:
        logger.error(f"[ai_core] Error extracting text from '{filepath}': {e}", exc_info=True)
        return None

def create_chunks_from_text(text: str, filename: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[Document]:
    if not text: return []
    try:
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        chunks_list = splitter.split_text(text)
        documents = [Document(page_content=chunk, metadata={"source": filename, "chunk_index": i}) for i, chunk in enumerate(chunks_list)]
        logger.info(f"[ai_core] Created {len(documents)} chunks for '{filename}'.")
        return documents
    except Exception as e:
        logger.error(f"[ai_core] Error creating chunks for '{filename}': {e}", exc_info=True)
        return []

def add_documents_to_vector_store(documents: list[Document]) -> bool:
    global vector_store
    if not embeddings:
        logger.error("[ai_core] Cannot add documents: Embeddings model not initialized.")
        return False
    if not documents:
        logger.info("[ai_core] No documents provided to add to vector store.")
        return True # Technically not a failure

    try:
        if vector_store is None: # If store doesn't exist or wasn't loaded
            logger.info("[ai_core] Creating new FAISS vector store from provided documents.")
            vector_store = FAISS.from_documents(documents, embeddings)
        else:
            logger.info(f"[ai_core] Adding {len(documents)} documents to existing FAISS vector store.")
            vector_store.add_documents(documents)
        
        save_vector_store() # Persist changes
        new_size = vector_store.index.ntotal if vector_store and vector_store.index else 0
        logger.info(f"[ai_core] Documents added/updated in vector store. New index size: {new_size}")
        return True
    except Exception as e:
        logger.error(f"[ai_core] Error adding documents to vector store: {e}", exc_info=True)
        return False

# --- RAG and Analysis ---
def generate_sub_queries(query: str) -> list[str]:
    if not llm:
        logger.error("[ai_core] LLM not initialized, cannot generate sub-queries for RAG.")
        return [query] # Fallback to original query
    try:
        # Ensure SUB_QUERY_PROMPT_TEMPLATE is a Langchain PromptTemplate instance
        if not isinstance(SUB_QUERY_PROMPT_TEMPLATE, PromptTemplate):
            logger.error("[ai_core] SUB_QUERY_PROMPT_TEMPLATE is not a valid PromptTemplate.")
            return [query]
            
        chain = LLMChain(llm=llm, prompt=SUB_QUERY_PROMPT_TEMPLATE)
        # The prompt expects 'query' and 'num_queries'
        response = chain.invoke({"query": query, "num_queries": MULTI_QUERY_COUNT})
        # Assuming chain.invoke returns a dict and the queries are in 'text' or similar key
        response_text = response.get('text', '') if isinstance(response, dict) else str(response)
        
        queries = [q.strip() for q in response_text.split('\n') if q.strip() and q.strip().lower() != "original query:"]
        logger.info(f"[ai_core] Generated sub-queries: {queries}")
        return queries if queries else [query] # Fallback if no sub-queries generated
    except Exception as e:
        logger.error(f"[ai_core] Error generating sub-queries: {e}", exc_info=True)
        return [query] # Fallback

def perform_rag_search(query: str) -> tuple[list[Document], str, dict]:
    if not vector_store:
        logger.warning("[ai_core] Vector store not available. RAG search cannot proceed.")
        return [], "Knowledge base is currently unavailable for search.", {}
    if not query:
        return [], "Cannot search with an empty query.", {}
    try:
        queries_to_search = generate_sub_queries(query) # This will include original query if sub-gen fails
        if query not in queries_to_search: # Ensure original query is always searched
            queries_to_search.insert(0, query)

        all_retrieved_docs_with_scores = []
        for q_text in set(queries_to_search): # Use set to avoid duplicate queries
            logger.debug(f"[ai_core] RAG searching for: '{q_text}' with k={RAG_SEARCH_K_PER_QUERY}")
            results_with_scores = vector_store.similarity_search_with_score(q_text, k=RAG_SEARCH_K_PER_QUERY)
            all_retrieved_docs_with_scores.extend(results_with_scores)
            logger.debug(f"[ai_core] Found {len(results_with_scores)} for sub-query '{q_text}'.")

        # Deduplicate documents and sort by score (lower is better for FAISS L2/IP)
        unique_docs_map = {} # {(source, content): (doc, score)}
        for doc_obj, score in all_retrieved_docs_with_scores:
            key = (doc_obj.metadata.get('source', 'Unknown'), doc_obj.page_content)
            if key not in unique_docs_map or score < unique_docs_map[key][1]:
                unique_docs_map[key] = (doc_obj, score)
        
        sorted_unique_docs_with_scores = sorted(unique_docs_map.values(), key=lambda item: item[1])
        
        # Take top RAG_CHUNK_K unique documents across all sub-query results
        final_docs_for_context = [doc_obj for doc_obj, score in sorted_unique_docs_with_scores[:RAG_CHUNK_K]]

        if not final_docs_for_context:
            logger.info(f"[ai_core] No relevant documents found after RAG search for query: '{query}'")
            return [], "No specific information found in the knowledge base for your query.", {}

        # Prepare context text and document map for citations
        context_parts = []
        docs_citation_map = {} # For simple [1], [2] style citations
        for i, doc in enumerate(final_docs_for_context):
            source_name = doc.metadata.get('source', f'Document {i+1}')
            chunk_idx_info = f" (chunk {doc.metadata.get('chunk_index', 'N/A')})" if 'chunk_index' in doc.metadata else ""
            context_parts.append(f"Source [{i+1}]: {source_name}{chunk_idx_info}\nContent: {doc.page_content}")
            docs_citation_map[i+1] = {"source": source_name, "content": doc.page_content} # Store full content for potential display

        context_text_for_llm = "\n\n---\n\n".join(context_parts)
        logger.info(f"[ai_core] RAG search complete. Returning {len(final_docs_for_context)} documents for LLM context. Citation map size: {len(docs_citation_map)}")
        
        return final_docs_for_context, context_text_for_llm, docs_citation_map
    except Exception as e:
        logger.error(f"[ai_core] Error during RAG search: {e}", exc_info=True)
        return [], "An error occurred while searching the knowledge base.", {}

def synthesize_chat_response(query: str, chat_history: str, tool_output: str = "", llm_preference: str = None) -> tuple[str, str]:
    """
    Synthesizes a chat response using the LLM, potentially in a ReAct style.
    The "Missing input keys: {"nodes", "content"}" error originates from how the LLM is invoked here.
    """
    if not llm: # Ensure main LLM is initialized
        logger.error("[ai_core] Main LLM (for chat) not initialized. Cannot synthesize response.")
        return "Error: The AI model for chat is not available at the moment.", "LLM not ready."
    
    # Select LLM if preference is provided and valid, otherwise use default
    current_llm_instance = llm # Default to the globally initialized llm
    if llm_preference and llm_preference != OLLAMA_MODEL: # Assuming OLLAMA_MODEL is the default
        try:
            logger.info(f"[ai_core] Attempting to use LLM preference: {llm_preference}")
            # You might have a dictionary of LLM instances or create on the fly
            # For simplicity, let's assume ChatOllama can switch model if KG_MODEL is also an Ollama model
            if "ollama" in llm_preference.lower(): # Basic check
                 alt_model_name = llm_preference.split('_', 1)[1] if '_' in llm_preference else llm_preference
                 current_llm_instance = ChatOllama(base_url=OLLAMA_BASE_URL, model=alt_model_name, request_timeout=OLLAMA_REQUEST_TIMEOUT)
                 current_llm_instance.invoke("test connection") # Quick test for the alternative LLM
                 logger.info(f"[ai_core] Switched to LLM: {alt_model_name}")
            else: # Fallback if preference is not an ollama model or format not recognized
                logger.warning(f"[ai_core] LLM preference '{llm_preference}' not fully supported or recognized as Ollama. Using default LLM.")
        except Exception as e_llm_pref:
            logger.error(f"[ai_core] Error switching to LLM preference '{llm_preference}': {e_llm_pref}. Using default LLM.")
            current_llm_instance = llm # Fallback to default

    current_thought = "Processing your request...\n"
    
    # --- CRITICAL: Adapt this `llm_input` to your SYNTHESIS_PROMPT_TEMPLATE ---
    # The error "Missing input keys: {"nodes", "content"}" means SYNTHESIS_PROMPT_TEMPLATE
    # (or the LLM call if not using a template directly) expects these keys.
    
    # Your SYNTHESIS_PROMPT_TEMPLATE is a PromptTemplate instance.
    # Find its `input_variables` attribute (e.g., print(SYNTHESIS_PROMPT_TEMPLATE.input_variables))
    # to see exactly what keys it expects.
    
    # Example: If SYNTHESIS_PROMPT_TEMPLATE.input_variables are ['query', 'chat_history', 'tool_output']
    # then this is correct:
    llm_input = {
        "query": query,                # This should map to the user's direct question
        "chat_history": chat_history,  # The conversational history
        "tool_output": tool_output     # Results from previous <ACTION> (e.g., RAG search results)
    }
    # If your prompt *also* specifically needs a key named "content" for the main query,
    # and "nodes" for retrieved documents (even if empty initially or part of tool_output),
    # you MUST add them here.
    # E.g., if template is: "Context (Nodes): {nodes}\nUser Question (Content): {content}\nHistory: {chat_history}..."
    # llm_input = {
    #     "content": query,
    #     "nodes": tool_output if "RAG_search Result:" in tool_output else "No specific documents were retrieved for this turn.",
    #     "chat_history": chat_history,
    #     "tool_observation": tool_output # Or maybe tool_output is split between nodes and observation
    # }
    # For now, assuming SYNTHESIS_PROMPT_TEMPLATE uses 'query', 'chat_history', 'tool_output'.

    current_thought += f"Prepared input for LLM based on SYNTHESIS_PROMPT_TEMPLATE. Keys expected: {SYNTHESIS_PROMPT_TEMPLATE.input_variables}\n"
    logger.debug(f"[ai_core] Input for LLM Chain (SYNTHESIS_PROMPT_TEMPLATE): {json.dumps(llm_input, indent=2)}")

    try:
        if not isinstance(SYNTHESIS_PROMPT_TEMPLATE, PromptTemplate):
            logger.error("[ai_core] SYNTHESIS_PROMPT_TEMPLATE is not a valid PromptTemplate instance!")
            return "Error: Server-side prompt configuration error.", "Prompt template invalid."

        chain = LLMChain(llm=current_llm_instance, prompt=SYNTHESIS_PROMPT_TEMPLATE)
        
        # LLMChain's run method can take keyword arguments matching the prompt's input_variables
        # or a single dictionary. Let's use keyword arguments for clarity if variables are simple.
        # If input_variables are complex, a dict is better.
        # response = chain.invoke(llm_input) # For LCEL chains
        # response_text = response.content if hasattr(response, 'content') else str(response)
        
        # For older LLMChain, .run() takes kwargs or a single dict.
        # Let's ensure all expected variables are passed.
        # Create a dictionary with only the keys expected by the prompt
        valid_llm_input = {key: llm_input.get(key) for key in SYNTHESIS_PROMPT_TEMPLATE.input_variables if key in llm_input}
        missing_keys = [key for key in SYNTHESIS_PROMPT_TEMPLATE.input_variables if key not in valid_llm_input]
        if missing_keys:
            logger.error(f"[ai_core] Critical Error: Missing keys for SYNTHESIS_PROMPT_TEMPLATE: {missing_keys}. Provided: {list(llm_input.keys())}")
            return f"Error: Internal prompt error. Missing expected inputs: {', '.join(missing_keys)}", "Prompt input construction failed."

        logger.info(f"[ai_core] Invoking LLM with model: {current_llm_instance.model if hasattr(current_llm_instance, 'model') else 'N/A'}")
        response_obj = chain.invoke(valid_llm_input) # chain.invoke expects a dict
        
        response_text = ""
        if isinstance(response_obj, dict) and 'text' in response_obj:
            response_text = response_obj['text']
        elif isinstance(response_obj, str):
            response_text = response_obj
        else: # Langchain LCEL might return AIMessage or similar
            response_text = str(response_obj.content) if hasattr(response_obj, 'content') else str(response_obj)


        current_thought += f"LLM raw response received.\n"
        logger.debug(f"[ai_core] LLM Raw Response: {response_text[:300]}...") # Log snippet

        # Simple parsing for <ACTION>...</ACTION> or direct speak
        # Your Node.js parsing logic seems more robust, this is simpler for Python side
        if "<ACTION>" in response_text and "</ACTION>" in response_text:
            # Extract action part
            action_content = response_text # Keep as is, let Node.js parse tools fully
            current_thought += "LLM proposed an action.\n"
            return action_content, current_thought.strip() # Return the full string with <ACTION> tags
        else:
            # Assume it's a direct answer if no <ACTION> tag
            current_thought += "LLM provided a direct answer.\n"
            return response_text.strip(), current_thought.strip()

    except Exception as e:
        logger.error(f"[ai_core] Error during LLM Chain invocation in synthesize_chat_response: {e}", exc_info=True)
        error_message = f"Error: Could not generate a response from the AI model. Details: {str(e)}"
        current_thought += f"LLM invocation failed: {str(e)}\n"
        return error_message, current_thought.strip()


def generate_document_analysis(filename: str, analysis_type: str, llm_preference: str = None) -> tuple[str | None, str]:
    global document_texts_cache
    if not llm: # Use the main llm instance
        logger.error("[ai_core] LLM not initialized, cannot generate document analysis.")
        return "Error: AI model is not available for analysis.", "LLM not ready."

    # LLM Selection for Analysis (similar to chat)
    current_llm_instance_analysis = llm
    if llm_preference and llm_preference != OLLAMA_MODEL:
        try:
            logger.info(f"[ai_core] Analysis using LLM preference: {llm_preference}")
            if "ollama" in llm_preference.lower():
                 alt_model_name_analysis = llm_preference.split('_', 1)[1] if '_' in llm_preference else llm_preference
                 current_llm_instance_analysis = ChatOllama(base_url=OLLAMA_BASE_URL, model=alt_model_name_analysis, request_timeout=OLLAMA_REQUEST_TIMEOUT)
                 current_llm_instance_analysis.invoke("test connection for analysis")
                 logger.info(f"[ai_core] Analysis switched to LLM: {alt_model_name_analysis}")
            else:
                logger.warning(f"[ai_core] Analysis LLM preference '{llm_preference}' not fully supported. Using default LLM.")
        except Exception as e_llm_pref_analysis:
            logger.error(f"[ai_core] Error switching for Analysis LLM preference '{llm_preference}': {e_llm_pref_analysis}. Using default.")
            current_llm_instance_analysis = llm


    try:
        text = document_texts_cache.get(filename)
        if not text:
            file_path_to_try = os.path.join(UPLOAD_FOLDER, filename)
            if not os.path.exists(file_path_to_try):
                file_path_to_try = os.path.join(DEFAULT_PDFS_FOLDER, filename)
            
            if os.path.exists(file_path_to_try):
                logger.info(f"[ai_core] Text for '{filename}' not in cache, loading from {file_path_to_try} for analysis.")
                text = extract_text_from_file(file_path_to_try)
                if text: document_texts_cache[filename] = text
            else:
                logger.error(f"[ai_core] File '{filename}' not found for analysis.")
                return None, "File not found for analysis."
        if not text: return None, "No text available in document for analysis."

        prompt_template_for_analysis = ANALYSIS_PROMPTS.get(analysis_type)
        if not prompt_template_for_analysis or not isinstance(prompt_template_for_analysis, PromptTemplate):
            logger.error(f"[ai_core] No valid prompt template found for analysis_type: '{analysis_type}'")
            return f"Error: Analysis type '{analysis_type}' is not configured correctly.", "Invalid analysis type."

        chain = LLMChain(llm=current_llm_instance_analysis, prompt=prompt_template_for_analysis)
        
        # Ensure the prompt's input variables are met. Usually 'doc_text_for_llm'.
        analysis_input = {"doc_text_for_llm": text[:ANALYSIS_MAX_CONTEXT_LENGTH]} # Truncate text
        logger.debug(f"[ai_core] Analysis input for LLM (type: {analysis_type}): {analysis_input['doc_text_for_llm'][:100]}...")

        response_obj = chain.invoke(analysis_input)
        response_text = ""
        if isinstance(response_obj, dict) and 'text' in response_obj:
            response_text = response_obj['text']
        elif isinstance(response_obj, str):
            response_text = response_obj
        else:
            response_text = str(response_obj.content) if hasattr(response_obj, 'content') else str(response_obj)

        # Assuming analysis doesn't have complex "thinking" part like ReAct, but could be added
        thinking_for_analysis = f"Analysis of type '{analysis_type}' performed on '{filename}'."
        logger.info(f"[ai_core] Analysis successful for '{filename}' ({analysis_type}).")
        return response_text.strip(), thinking_for_analysis
        
    except Exception as e:
        logger.error(f"[ai_core] Error in document analysis for '{filename}' ({analysis_type}): {e}", exc_info=True)
        return f"Error performing analysis: {str(e)}", f"Analysis failed due to: {str(e)}"

# --- Knowledge Graph Generation Functions ---
# (Your KG functions: get_kg_filepath, load_knowledge_graph, _kg_split_into_chunks,
#  _kg_process_single_chunk, _kg_merge_graphs, _kg_save_graph, generate_knowledge_graph_from_pdf
#  Seem largely okay from what you provided, ensure KG_MODEL, KG_PROMPT_TEMPLATE are correctly configured
#  and _initialize_kg_ollama_client works.)
#  Ensure KG_PROMPT_TEMPLATE is a string, not a PromptTemplate object if _kg_process_single_chunk formats it.
#  If KG_PROMPT_TEMPLATE is a PromptTemplate object, then _kg_process_single_chunk needs to use LLMChain.

# Make sure all KG functions are defined as in your provided `ai_core.py`
# For brevity, I'm not re-listing them here, but they should be present.
# Key function:
# def generate_knowledge_graph_from_pdf(doc_filename: str) -> dict | None:
#    ... uses _initialize_kg_ollama_client, _kg_split_into_chunks, _kg_process_single_chunk, _kg_merge_graphs, _kg_save_graph ...
#    ... ensure KG_PROMPT_TEMPLATE is correctly used with ollama.Client.chat ...
#
# def format_kg_info_for_llm(kg_data: dict, query: str, max_insights: int = 3, max_details_per_node: int = 2) -> str:
#    ... your existing logic ...


# --- Ensure all KG helper functions from your ai_core.py are here ---
def get_kg_filepath(doc_filename: str) -> str:
    base_name = os.path.splitext(doc_filename)[0]
    kg_file = f"{base_name}{KG_FILENAME_SUFFIX}"
    return os.path.join(KG_OUTPUT_FOLDER, kg_file)

def load_knowledge_graph(doc_filename: str) -> dict | None:
    kg_filepath = get_kg_filepath(doc_filename)
    if os.path.exists(kg_filepath):
        try:
            with open(kg_filepath, 'r', encoding='utf-8') as f:
                kg_data = json.load(f)
            logger.info(f"[ai_core] Loaded KG for '{doc_filename}' from {kg_filepath}")
            return kg_data
        except Exception as e:
            logger.error(f"[ai_core] Error loading KG JSON from {kg_filepath}: {e}", exc_info=True)
    else:
        logger.info(f"[ai_core] No KG found for '{doc_filename}' at {kg_filepath}") # Changed to info
    return None

def _kg_split_into_chunks(text: str, chunk_size: int, overlap: int) -> list[str]:
    # (Your existing logic from provided ai_core.py)
    logger.info(f"[ai_core KG] Splitting text into chunks (size={chunk_size}, overlap={overlap})...")
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        next_start = end - overlap
        if next_start <= start and end < text_len : start = start + 1 
        else: start = next_start if end < text_len else end
    logger.info(f"[ai_core KG] Split into {len(chunks)} chunks.")
    return chunks


def _kg_process_single_chunk(chunk_data: tuple[int, str], ollama_client_instance: ollama.Client, model_name: str, prompt_template_str: str) -> dict | None:
    # (Your existing logic, ensure prompt_template_str is formatted correctly if it's a string)
    index, chunk_text = chunk_data; chunk_num = index + 1
    full_prompt = prompt_template_str.format(chunk_text=chunk_text) 
    if not ollama_client_instance: logger.error(f"[ai_core KG] Ollama client NA for chunk {chunk_num}."); return None
    try:
        response = ollama_client_instance.chat(model=model_name, messages=[{"role": "user", "content": full_prompt}], format="json", options={"num_ctx": 4096, "temperature": 0.3})
        content = response.get('message', {}).get('content', '')
        if not content: logger.warning(f"[ai_core KG] Empty response chunk {chunk_num}"); return {}
        if content.strip().startswith('```json'): content = content.strip()[7:-3].strip()
        elif content.strip().startswith('```'): content = content.strip()[3:-3].strip()
        graph_data = json.loads(content)
        if isinstance(graph_data, dict) and 'nodes' in graph_data and isinstance(graph_data['nodes'], list) and 'edges' in graph_data and isinstance(graph_data['edges'], list): return graph_data
        logger.warning(f"[ai_core KG] Invalid graph structure chunk {chunk_num}. Content: {content[:200]}..."); return {}
    except json.JSONDecodeError as je: logger.error(f"[ai_core KG] JSON Decode Error chunk {chunk_num}: {je}. Content: {content[:500]}..."); return {}
    except Exception as e: logger.error(f"[ai_core KG] Error processing chunk {chunk_num}: {e}", exc_info=True); return {}

def _kg_merge_graphs(graphs: list[dict]) -> dict:
    # (Your existing complex merging logic from provided ai_core.py)
    logger.info("[ai_core KG] Merging graph fragments...")
    final_nodes = {}; final_edges = set()
    for i, gf in enumerate(graphs):
        if gf is None or not (isinstance(gf, dict) and 'nodes' in gf and 'edges' in gf): logger.warning(f"[ai_core KG] Skipping invalid fragment {i}"); continue
        for node in gf.get('nodes', []):
            if not isinstance(node, dict): logger.warning(f"[ai_core KG] Skipping non-dict node {node}"); continue
            node_id = node.get('id'); 
            if not (node_id and isinstance(node_id, str)): logger.warning(f"[ai_core KG] Skipping node invalid ID {node}"); continue
            if node_id not in final_nodes: final_nodes[node_id] = node
            else: # Merge properties
                en = final_nodes[node_id]
                for k in ['description', 'type', 'parent']: 
                    if node.get(k) and (not en.get(k) or len(str(node.get(k))) > len(str(en.get(k)))): 
                        if k == 'parent' and node.get(k) == node_id: continue
                        en[k] = node.get(k)
        for edge in gf.get('edges', []):
            if not (isinstance(edge, dict) and all(k in edge for k in ['from', 'to', 'relationship'])): logger.warning(f"[ai_core KG] Skipping invalid edge {edge}"); continue
            if not all(isinstance(edge.get(k), str) and edge.get(k) for k in ['from', 'to', 'relationship']): logger.warning(f"[ai_core KG] Skipping edge non-string components {edge}"); continue
            if edge['from'] == edge['to']: continue # Optional: skip self-loops
            final_edges.add((edge['from'], edge['to'], edge['relationship']))
    mgn = list(final_nodes.values()); vni = set(n['id'] for n in mgn)
    fme = [{"from": e[0], "to": e[1], "relationship": e[2]} for e in final_edges if e[0] in vni and e[1] in vni]
    mg = {"nodes": mgn, "edges": fme}; logger.info(f"[ai_core KG] Merged {len(mg['nodes'])} nodes, {len(mg['edges'])} edges."); return mg


def _kg_save_graph(graph: dict, original_doc_filename: str):
    kg_filepath = get_kg_filepath(original_doc_filename)
    try:
        logger.info(f"[ai_core KG] Saving graph for '{original_doc_filename}' to {kg_filepath}...")
        os.makedirs(os.path.dirname(kg_filepath), exist_ok=True)
        with open(kg_filepath, 'w', encoding='utf-8') as f: json.dump(graph, f, indent=2, ensure_ascii=False)
        logger.info(f"[ai_core KG] Graph for '{original_doc_filename}' saved to {kg_filepath}.")
    except Exception as e: logger.error(f"Failed to save graph to {kg_filepath}: {e}", exc_info=True)


def generate_knowledge_graph_from_pdf(doc_filename: str) -> dict | None:
    global document_texts_cache
    kg_ollama_client = _initialize_kg_ollama_client()
    if not kg_ollama_client: logger.error(f"[ai_core KG] Cannot gen graph for '{doc_filename}'. KG Ollama client NA."); return None
    logger.info(f"[ai_core KG] Generating KG for '{doc_filename}'...")
    doc_text = document_texts_cache.get(doc_filename)
    if not doc_text:
        logger.debug(f"[ai_core KG] '{doc_filename}' not cached. Loading from disk...")
        fp_load = os.path.join(UPLOAD_FOLDER, doc_filename)
        if not os.path.exists(fp_load): fp_load = os.path.join(DEFAULT_PDFS_FOLDER, doc_filename)
        if os.path.exists(fp_load):
            logger.info(f"[ai_core KG] Found '{doc_filename}' at {fp_load}. Extracting...")
            doc_text = extract_text_from_file(fp_load)
            if doc_text: document_texts_cache[doc_filename] = doc_text; logger.info(f"[ai_core KG] Cached text for '{doc_filename}'.")
            else: logger.error(f"[ai_core KG] Failed to extract text from '{fp_load}'."); return None
        else: logger.error(f"[ai_core KG] File '{doc_filename}' not found."); return None
    if not doc_text.strip(): logger.error(f"[ai_core KG] No text content for '{doc_filename}'."); return None
    chunks = _kg_split_into_chunks(doc_text, KG_CHUNK_SIZE, KG_CHUNK_OVERLAP)
    if not chunks: logger.warning(f"[ai_core KG] No chunks for '{doc_filename}'. Aborted."); return None
    logger.info(f"[ai_core KG] Processing {len(chunks)} chunks for '{doc_filename}' using {KG_MAX_WORKERS} workers, model {KG_MODEL}.")
    all_partial_graphs = []
    prompt_str = KG_PROMPT_TEMPLATE # KG_PROMPT_TEMPLATE should be a string from config
    with concurrent.futures.ThreadPoolExecutor(max_workers=KG_MAX_WORKERS) as executor:
        tasks = [(i, chunk) for i, chunk in enumerate(chunks)]
        future_map = {executor.submit(_kg_process_single_chunk, task, kg_ollama_client, KG_MODEL, prompt_str): task[0] for task in tasks}
        for future in tqdm(concurrent.futures.as_completed(future_map), total=len(tasks), desc=f"KG Chunks ({doc_filename})"):
            try: frag = future.result();
            except Exception as e_f: logger.error(f"[ai_core KG] Error from chunk {future_map[future]+1} for '{doc_filename}': {e_f}", exc_info=True); continue
            if frag: all_partial_graphs.append(frag)
    logger.info(f"[ai_core KG] Processed {len(all_partial_graphs)} valid fragments from {len(chunks)} chunks for '{doc_filename}'.")
    if not all_partial_graphs: logger.error(f"[ai_core KG] No valid graph fragments for '{doc_filename}'."); return None
    final_graph = _kg_merge_graphs(all_partial_graphs)
    _kg_save_graph(final_graph, doc_filename)
    logger.info(f"[ai_core KG] KG generation for '{doc_filename}' complete.")
    return final_graph

# --- End KG helper functions ---

if __name__ == "__main__":
    # (Your existing __main__ test block from provided ai_core.py)
    logging.basicConfig(level=logging.DEBUG, format="%(asctime)s - %(levelname)s - [%(name)s:%(lineno)d] - %(message)s")
    for dir_path in [DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, FAISS_FOLDER, KG_OUTPUT_FOLDER]: os.makedirs(dir_path, exist_ok=True)
    logger.info("AI Core Testing Main...")
    embeddings_instance, llm_instance = initialize_ai_components()
    if not (embeddings_instance and llm_instance): logger.error("Failed to init AI. Aborting."); exit()
    load_vector_store()
    # ... (rest of your test logic for dummy PDF, RAG, KG, and chat) ...
    # Example simplified test:
    test_pdf = "DevOps-Test-KG.pdf" # Make sure this file exists in DEFAULT_PDFS_FOLDER or UPLOAD_FOLDER
    if os.path.exists(os.path.join(DEFAULT_PDFS_FOLDER, test_pdf)) or os.path.exists(os.path.join(UPLOAD_FOLDER, test_pdf)):
        generate_knowledge_graph_from_pdf(test_pdf)
        test_query = "What is DevOps?"
        answer, thinking = synthesize_chat_response(query=test_query, chat_history="", tool_output="")
        logger.info(f"\nQuery: {test_query}\nThinking: {thinking}\nAnswer: {answer}")
    else:
        logger.warning(f"Test PDF '{test_pdf}' not found. Skipping some tests.")

    logger.info("\nAI Core test completed.")

# --- END OF FILE ai_core.py ---